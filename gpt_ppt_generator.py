import openai
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import time

from dotenv import load_dotenv

# Load .env file
load_dotenv()

# Retrieve API key from environment variables
openai.api_key = os.getenv("OPENAI_API_KEY")

template_path = 'templates/template.pptx'
output_directory = 'ppt_files'
slide_notes = []
launguage = "English"
step_by_step = False

ppt_title = "cyber fusion center"
additional_info = "1. what is cyber fusion center and why is called cyber fusion center; 2. key functions and responsibilities of cyber fusion center; 3. the staffing of cyber fusion center; 4. the management of cyber fusion center; 5. the requirements of the role of head of cyber fusion center"

#ppt_title = "cyber fusion center"
#additional_info = "1. what is cyber fusion center and why is called cyber fusion center; 2. key functions and responsibilities of cyber fusion center; 3. the staffing of cyber fusion center; 4. the management of cyber fusion center; 5. the requirements of the role of head of cyber fusion center"


#ppt_title = "NIST Cybersecurity Framework"
#additional_info = "1. what is NIST Cybersecurity Framework; 2. how it works; 3. how NIST Cybersecurity Framework can enterprises to improve security. 4.what the the differences between NIST Cybersecurity Framework and other cybersecurity framework like MITRE ATT&CK Framework. "

#ppt_title = "MITRE ATT&CK Framework"
#additional_info = "1. what is MITRE ATT&CK Framework; 2. how it works; 3. how MITRE ATT&CK Framework can enterprises to improve security. "

#ppt_title = "device fingerprinting and device risk identification"
#additional_info = "1. what id device fingerprinting; 2. how it works; 3. how device fingerprinting can help identify device risk 4. key use cases; 5. key players in the market. "
#additional_info = "1. the diffences between tranditional firewall and cloud firewall; 2. key features of cloud firwell; 3. the comparison of key players in cloud firewall market; 4. the future of cloud firewall market. "


def generate_pptx(outlines, ppt_title, template_path, output_directory, add_notes=False):
    # Create a new PowerPoint presentation from template
    # and fill it with generated key points (slides) and text.
    # See 'pptx' library documentation for further graphic customization.
    prs = Presentation(template_path)
    # Add a new slide with a title and content
    slide_layout = prs.slide_layouts[0]  # Use the second layout in the template
    slide = prs.slides.add_slide(slide_layout)   
    # Add a title to the slide
    title_placeholder = slide.shapes.title
    title_placeholder.text = ppt_title    
         
    for i, outline in enumerate(outlines):
        lines=outline.split('\n')
        topic=lines[0]
        text='\n'.join(lines[1:])
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = topic
        title_shape = slide.shapes.title
        title_shape.text = topic

        # Format title font.
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.size = Pt(42)
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(9)
        height = Inches(4)

        # Add content to the slide
        content_placeholder = slide.placeholders[1]  # The first placeholder is the title
        content_placeholder.text = text
        

        if add_notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_notes[i]
            
    # Save presentation.
    pptx_filename = os.path.join(output_directory, f"{ppt_title}-{launguage}.pptx")
    prs.save(pptx_filename)     

# Define a function to ensure the conversation history does not exceed the maximum token limit
def ensure_under_token_limit(conversation, max_tokens):
    total_tokens = sum([len(message['content'].split()) for message in conversation])
    
    while total_tokens > max_tokens:
        # Remove the earliest messages until the total token count of the conversation history does not exceed the maximum token limit
        removed_message = conversation.pop(0)
        total_tokens -= len(removed_message['content'].split())
    
    return conversation

# Start a new conversation
init_conversation = [
    {"role": "system", "content": "You are a helpful assistant."},
    {"role": "user", "content": f"Please draft a presentation outlines about {ppt_title} in {launguage} for a 20-minutes talk; in the presentation, please mention following key points: {additional_info}."}
]

# Step 1: Generate the outline for the PowerPoint
outline_response = openai.ChatCompletion.create(
  model="gpt-3.5-turbo",
  messages=init_conversation
)

outline = outline_response['choices'][0]['message']['content']


# Add the model's response to the conversation history
init_conversation.append({"role": "assistant", "content": f"{outline}"})

# Extract each page's topic from the outline
# Here it is assumed that each topic is on a separate line
topics = outline.split('\n')
topics = topics[2:-1]

# First, combine all topics into a string with '\n' as the separator
topics_str = '\n'.join(topics)

# Then, split the string into multiple parts with empty lines (i.e., two consecutive '\n') as the separator
page_content = topics_str.split('\n\n')

# Finally, remove the leading and trailing white space (including '\n') from each part

outlines = [content.strip() for content in page_content]

print(f"presentation outline: {outlines}")

if step_by_step:
    # Generate PowerPoint with the outlines
    generate_pptx(outlines, ppt_title, template_path, output_directory)
    choice = input("Press 'C' to call OpenAI to generate presentation notes, or any other key to exit: ")
    if choice != 'C' and choice != 'c':
        sys.exit()

conversation = init_conversation
# Open the file to write the content
with open(f'ppt_files/{ppt_title}_ppt_notes-{launguage}.txt', 'w', encoding='utf8') as f:   
# Step 2: Generate content for each topic
    for i, topic in enumerate(outlines):
        #conversation = init_conversation
        conversation.append({"role": "user", "content": f"draft the presenation notes in {launguage} in a spoken style to elaborate slide {i}: {topic}"})

        # Ensure the conversation history does not exceed the maximum token limit
        conversation = ensure_under_token_limit(conversation, 3200)

        presentation_notes_response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=conversation
        )
        time.sleep(1)

        # Add the model's response to the conversation history
        conversation.append({"role": "assistant", "content": presentation_notes_response['choices'][0]['message']['content']})

        presentation_notes = presentation_notes_response['choices'][0]['message']['content']
        slide_notes.append(presentation_notes)
        # Write the content to the file
        f.write(f"Slide {i+1} - {topic}: {presentation_notes}\n")        
        print(f"Slide {i+1} - {topic}: {presentation_notes}")

        

if step_by_step:
    choice = input("Press 'Y' to regenerate presentation ppt file with notes, or any other key to exit: ")
    if choice != 'Y' and choice != 'y':
        sys.exit()
    # Generate PowerPoint using outlines
generate_pptx(outlines, ppt_title, template_path, output_directory, add_notes=True)




